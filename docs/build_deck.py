# -*- coding: utf-8 -*-
"""
عرض تنفيذي مبسّط لتطوير منظومة المشتريات — مجموعة الذيابي للمقاولات.
خط Neo Sans Arabic، جداول من عمودين (الوضع ← المعالجة) بعبارات قصيرة،
مقسّم إلى ثلاثة أجزاء واضحة. عروض الأسعار الثلاثة إلزامية لأي شراء؛
وسجل الأسعار لتسعير المناقصات فقط.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

INK = RGBColor(0x0E, 0x2A, 0x38)
INK2 = RGBColor(0x14, 0x3A, 0x4C)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
SAND = RGBColor(0xF6, 0xF3, 0xED)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xDD, 0xD4, 0xC2)
SLATE = RGBColor(0x33, 0x40, 0x4A)
MUTE = RGBColor(0x6B, 0x74, 0x7C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0x9E, 0x2B, 0x25)
AMBER = RGBColor(0xC5, 0x7A, 0x0C)
GREEN = RGBColor(0x2E, 0x6B, 0x3E)
ROWALT = RGBColor(0xF1, 0xEC, 0xE2)

F = "Neo Sans Arabic"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def rtl(p):
    p._p.get_or_add_pPr().set('rtl', '1')


def table_rtl(tb):
    """تفعيل اتجاه RTL للجدول على عنصر tblPr (المكان الصحيح)."""
    from pptx.oxml.ns import qn
    tbl = tb._tbl
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is None:
        tblPr = tbl.makeelement(qn('a:tblPr'), {})
        tbl.insert(0, tblPr)
    tblPr.set('rtl', '1')


def rect(s, l, t, w, h, fill, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp


def txt(s, l, t, w, h, lines, *, size=16, color=SLATE, bold=False,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, is_rtl=True, space=6):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(6); tf.margin_top = tf.margin_bottom = Pt(3)
    if isinstance(lines, str):
        lines = [lines]
    for i, it in enumerate(lines):
        d = it if isinstance(it, dict) else {'t': it}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if is_rtl:
            rtl(p)
        p.alignment = d.get('align', align); p.space_after = Pt(d.get('space', space))
        r = p.add_run(); r.text = d['t']
        r.font.size = Pt(d.get('size', size)); r.font.bold = d.get('bold', bold)
        r.font.name = F; r.font.color.rgb = d.get('color', color)
    return tb


def bullets(s, l, t, w, h, items, size=16, gap=10, color=SLATE, mark="•  "):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        d = it if isinstance(it, dict) else {'t': it}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rtl(p); p.alignment = PP_ALIGN.RIGHT; p.space_after = Pt(gap)
        r = p.add_run(); r.text = mark + d['t']
        r.font.size = Pt(d.get('size', size)); r.font.bold = d.get('bold', False)
        r.font.name = F; r.font.color.rgb = d.get('color', color)
    return tb


def page(s, idx=None):
    txt(s, Inches(0.3), SH - Inches(0.46), Inches(2), Inches(0.32),
        f"{len(slides)} / {TOTAL}", size=10, color=MUTE, align=PP_ALIGN.LEFT, is_rtl=False)


def brandbar(s, title, part):
    rect(s, 0, 0, SW, SH, SAND)
    rect(s, 0, 0, SW, Inches(1.12), INK)
    rect(s, 0, Inches(1.12), SW, Pt(4), GOLD)
    txt(s, Inches(0.35), Inches(0.2), Inches(9.6), Inches(0.62), title, size=25,
        color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    txt(s, Inches(0.35), Inches(0.78), Inches(9.6), Inches(0.32), part, size=13,
        color=GOLD, bold=True, align=PP_ALIGN.RIGHT)
    # شعار يسار
    txt(s, Inches(0.35), Inches(0.3), Inches(3.2), Inches(0.55),
        'مجموعة الذيابي', size=14, color=GOLD, bold=True, align=PP_ALIGN.LEFT)


def decision(s, text):
    y = SH - Inches(1.0)
    rect(s, Inches(0.35), y, SW - Inches(0.7), Inches(0.74), INK2)
    rect(s, SW - Inches(0.35) - Pt(6), y, Pt(6), Inches(0.74), GOLD)
    txt(s, Inches(0.55), y + Inches(0.04), SW - Inches(1.1), Inches(0.66),
        [{'t': 'المطلوب', 'size': 11, 'color': GOLD, 'bold': True, 'space': 2},
         {'t': text, 'size': 14, 'color': WHITE}],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


def two_col(s, top, rows, *, fsize=15, row_h=0.74):
    """جدول مبسّط: الوضع الحالي (يمين) ← المعالجة المقترحة (يسار)."""
    l = Inches(0.5); w = SW - Inches(1.0)
    g = s.shapes.add_table(len(rows) + 1, 2, l, top, w, Inches(row_h * (len(rows) + 1)))
    tb = g.table; table_rtl(tb)
    tb.columns[0].width = Emu(int(w * 0.5)); tb.columns[1].width = Emu(int(w * 0.5))
    for j, (h, bg) in enumerate([('الوضع الحالي', INK), ('المعالجة المقترحة', GREEN)]):
        c = tb.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = bg
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Pt(6); c.margin_top = c.margin_bottom = Pt(3)
        p = c.text_frame.paragraphs[0]; rtl(p); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = F
    for i, (prob, sol) in enumerate(rows, start=1):
        for j, (text, col, bd) in enumerate([(prob, INK, True), (sol, GREEN, False)]):
            c = tb.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 else ROWALT
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Pt(8); c.margin_top = c.margin_bottom = Pt(4)
            p = c.text_frame.paragraphs[0]; rtl(p); p.alignment = PP_ALIGN.RIGHT
            r = p.add_run(); r.text = text; r.font.size = Pt(fsize); r.font.name = F
            r.font.color.rgb = col; r.font.bold = bd
    return g


def proc_col(s, l, top, w, rows, row_h=0.7):
    """عمود إجراءات: المرحلة + الإجراء + المسؤول."""
    g = s.shapes.add_table(len(rows) + 1, 3, l, top, w, Inches(row_h * (len(rows) + 1)))
    tb = g.table; table_rtl(tb)
    cw = (3.0, 4.0, 2.6); total = sum(cw)
    for i, c in enumerate(cw):
        tb.columns[i].width = Emu(int(w * c / total))
    for j, h in enumerate(('المرحلة', 'الإجراء', 'المسؤول')):
        c = tb.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = INK
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Pt(4); c.margin_top = c.margin_bottom = Pt(2)
        p = c.text_frame.paragraphs[0]; rtl(p); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h; r.font.size = Pt(12.5); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = F
    cols = [(INK, True, PP_ALIGN.RIGHT), (SLATE, False, PP_ALIGN.RIGHT), (INK, False, PP_ALIGN.CENTER)]
    for i, row in enumerate(rows, start=1):
        for j, cell in enumerate(row):
            c = tb.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 else ROWALT
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Pt(5); c.margin_top = c.margin_bottom = Pt(2)
            col, bd, al = cols[j]
            p = c.text_frame.paragraphs[0]; rtl(p); p.alignment = al
            r = p.add_run(); r.text = cell; r.font.size = Pt(12.5); r.font.name = F
            r.font.color.rgb = col; r.font.bold = bd
    return g


def kpi_tile(s, x, y, w, h, big, label, co, sub=None):
    rect(s, x, y, w, h, CARD, line=LINE)
    rect(s, x, y, w, Pt(7), co)
    txt(s, x + Inches(0.05), y + Inches(0.14), w - Inches(0.1), Inches(0.62),
        big, size=34, color=co, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.05), y + h - Inches(0.62), w - Inches(0.1), Inches(0.56),
        ([{'t': label, 'size': 13.5, 'color': INK, 'bold': True, 'space': 1}]
         + ([{'t': sub, 'size': 11, 'color': MUTE}] if sub else [])),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def hbars(s, l, t, w, rows, maxv, lblw=Inches(2.2), bar_h=Inches(0.42), gap=Inches(0.22)):
    """أعمدة أفقية: التسمية يميناً، والعمود يمتدّ نحو اليسار."""
    track_r = l + w - lblw - Inches(0.15)
    track_w = w - lblw - Inches(0.15)
    for i, (lbl, val, co) in enumerate(rows):
        y = t + i * (bar_h + gap)
        txt(s, l + w - lblw, y - Inches(0.02), lblw, bar_h + Inches(0.04), lbl,
            size=14, color=INK, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, l, y, track_w, bar_h, RGBColor(0xEC, 0xE6, 0xDA))
        bw = max(int(track_w * (val / maxv)), Inches(0.5))
        rect(s, track_r - bw, y, bw, bar_h, co)
        txt(s, track_r - bw + Inches(0.05), y, bw - Inches(0.1), bar_h, str(val),
            size=14, color=WHITE, bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
            is_rtl=False)


def source_note(s, text='المصدر: نظام المشتريات — سجل أوامر الشراء (53 أمراً)'):
    txt(s, Inches(0.35), Inches(0.66), Inches(5.0), Inches(0.4), text,
        size=10.5, color=SAND, align=PP_ALIGN.LEFT, is_rtl=True)


# أرقام سجل الأسعار (من نظام المشتريات الفعلي)
PR_ITEMS = '585'    # عدد الأصناف المسجّلة (579 مُسعّر، 6 بدون سعر)
PR_PRICES = '864'   # عدد السجلات السعرية

TOTAL = 20
slides = []
P1 = 'الجزء الأول · المشاكل ومعالجتها'
P2 = 'الجزء الثاني · المنظومة المقترحة'
P3 = 'الجزء الثالث · التنفيذ والقرار'


def new():
    s = prs.slides.add_slide(BLANK); slides.append(s); return s


# 1) الغلاف
s = new()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(3.05), SW, Pt(4), GOLD)
txt(s, Inches(1), Inches(0.95), SW - Inches(2), Inches(0.5),
    'مجموعة الذيابي للمقاولات', size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(1.95), SW - Inches(1.6), Inches(1.0),
    'تطوير منظومة المشتريات', size=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.25), SW - Inches(2), Inches(0.6),
    'المشاكل ومعالجتها · المنظومة المقترحة · المطلوب من الإدارة', size=18,
    color=GOLD, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.7), SW - Inches(2), Inches(1.2),
    [{'t': 'مرفوع إلى سعادة المدير العام والإدارة العليا', 'size': 16, 'color': WHITE, 'bold': True, 'space': 8},
     {'t': 'إعداد: إدارة العمليات والمشتريات', 'size': 14, 'color': SAND}],
    align=PP_ALIGN.CENTER)

# 2) المحتوى
s = new(); brandbar(s, 'محتوى العرض', 'ثلاثة أجزاء')
agenda = [
    ('1', 'الوضع والمشاكل', 'أرقام النظام، ثم سبعة محاور بمعالجاتها', GREEN),
    ('2', 'المنظومة المقترحة', 'السياسات · الدورة الكاملة · مصفوفة الصلاحيات', AMBER),
    ('3', 'التنفيذ والقرار', 'خارطة الطريق · مؤشرات النجاح · القرارات المطلوبة', INK),
]
for i, (n, ti, de, co) in enumerate(agenda):
    y = Inches(1.7) + i * Inches(1.5)
    rect(s, Inches(0.6), y, SW - Inches(1.2), Inches(1.2), CARD, line=LINE)
    rect(s, SW - Inches(1.85), y, Inches(1.25), Inches(1.2), co)
    txt(s, SW - Inches(1.85), y, Inches(1.25), Inches(1.2), n, size=46, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, is_rtl=False)
    txt(s, Inches(0.8), y + Inches(0.12), SW - Inches(2.9), Inches(1.0),
        [{'t': ti, 'size': 24, 'color': INK, 'bold': True, 'space': 4},
         {'t': de, 'size': 15, 'color': SLATE}], align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE)
page(s, 2)

# 3) الملخص التنفيذي
s = new(); brandbar(s, 'الملخص التنفيذي', 'الفكرة في أربع نقاط')
cards = [
    ('المشكلة', 'لا توجد قناة موحّدة ولا دورة مُلزِمة ولا صلاحيات واضحة للمشتريات.', RED),
    ('الأثر', 'تأخّر 68٪ من الأوامر، وازدواج إنفاق، وتآكل ثقة الموردين.', AMBER),
    ('الحل', 'قناة واحدة + دورة موثّقة + ثلاثة عروض لكل شراء + صلاحيات واضحة.', GREEN),
    ('المطلوب', 'إقرار السياسات وخارطة الطريق من الإدارة العليا.', INK),
]
cw = (SW - Inches(1.2)) / 2
for i, (ti, de, co) in enumerate(cards):
    r, c = divmod(i, 2)
    x = SW - Inches(0.5) - cw - c * (cw + Inches(0.2)); y = Inches(1.6) + r * Inches(1.7)
    rect(s, x, y, cw, Inches(1.5), CARD, line=LINE)
    rect(s, x + cw - Pt(9), y, Pt(9), Inches(1.5), co)
    txt(s, x + Inches(0.2), y + Inches(0.18), cw - Inches(0.5), Inches(1.2),
        [{'t': ti, 'size': 22, 'color': co, 'bold': True, 'space': 6},
         {'t': de, 'size': 16, 'color': SLATE}], align=PP_ALIGN.RIGHT)
page(s, 3)

# 4) لوحة التحديات
s = new(); brandbar(s, 'أبرز التحديات', 'مرتّبة حسب الأهمية')
rowsT = [
    ('ازدواجية التوريد', 'حرج', RED), ('الشراء المباشر خارج المشتريات', 'حرج', RED),
    ('بطء التحويل المالي', 'حرج', RED), ('الإفراط في صفة «عاجل»', 'مرتفع', AMBER),
    ('الطلبات المعلّقة', 'مرتفع', AMBER), ('غياب الدورة والصلاحيات', 'مرتفع', AMBER),
    ('بطء تسعير المناقصات', 'متوسط', GREEN), ('عهدة قصر الياسمين', 'متوسط', GREEN),
]
for i, (t, lvl, co) in enumerate(rowsT):
    c, r = divmod(i, 4)
    colw = (SW - Inches(1.1)) / 2
    x = Inches(0.5) + (1 - c) * (colw + Inches(0.1)); y = Inches(1.7) + r * Inches(1.1)
    rect(s, x, y, colw, Inches(0.92), CARD, line=LINE)
    rect(s, x + colw - Inches(1.15), y, Inches(1.15), Inches(0.92), co)
    txt(s, x + colw - Inches(1.15), y, Inches(1.15), Inches(0.92), lvl, size=15,
        color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.15), y, colw - Inches(1.35), Inches(0.92), t, size=16,
        color=INK, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
page(s, 4)

# لوحة 1) الوضع الراهن بالأرقام
s = new(); brandbar(s, 'الوضع الراهن بالأرقام', 'من واقع نظام المشتريات')
tiles = [
    ('57', 'أمر شراء', INK, 'في النظام'),
    ('53', 'مورد معتمد', INK, 'شركاء نشطون'),
    ('68٪', 'طلبات متأخرة', RED, 'من تحليل سجل الأوامر'),
    ('19', 'يوم متوسط التأخير', RED, 'أقصى تأخير 52 يوماً'),
    ('23٪', 'لم تكتمل', AMBER, 'ملغاة أو قيد المراجعة'),
    ('405', 'ألف ريال', GOLD, 'إجمالي قيمة الأوامر'),
]
cw = (SW - Inches(1.2)) / 3
for i, (big, lbl, co, sub) in enumerate(tiles):
    r, c = divmod(i, 3)
    x = SW - Inches(0.5) - (c + 1) * cw - c * Inches(0.1) + Inches(0.1)
    y = Inches(1.55) + r * Inches(2.05)
    kpi_tile(s, x, y, cw, Inches(1.85), big, lbl, co, sub)
source_note(s)
page(s)

# لوحة 2) حالة الأوامر والأولوية  (الحالة يميناً، الأولوية يساراً)
s = new(); brandbar(s, 'لوحة حالة الأوامر والأولوية', 'توزيع أوامر الشراء')
RX = Inches(6.95); RW = Inches(5.9)   # النصف الأيمن
LX = Inches(0.5); LW = Inches(5.9)    # النصف الأيسر
txt(s, RX, Inches(1.35), RW, Inches(0.4), 'حسب الحالة', size=16, color=GOLD,
    bold=True, align=PP_ALIGN.RIGHT)
hbars(s, RX, Inches(1.95), RW, [
    ('مُسلَّم', 38, GREEN), ('قيد المراجعة', 7, AMBER), ('ملغى', 5, RED),
    ('مُسلَّم جزئياً', 2, AMBER), ('قيد التوريد', 1, INK),
], maxv=38, lblw=Inches(1.9))
txt(s, LX, Inches(1.35), LW, Inches(0.4), 'حسب الأولوية', size=16, color=GOLD,
    bold=True, align=PP_ALIGN.RIGHT)
hbars(s, LX, Inches(1.95), LW, [
    ('متوسط', 22, INK), ('عالي', 18, AMBER), ('عاجل', 13, RED),
], maxv=22, lblw=Inches(1.6))
# تنبيه العاجل (أسفل عمود الأولوية يساراً)
rect(s, LX, Inches(4.0), LW, Inches(1.35), RGBColor(0xF7, 0xE9, 0xE7), line=RED, lw=1.5)
txt(s, LX + Inches(0.2), Inches(4.15), LW - Inches(0.4), Inches(1.1),
    [{'t': 'مفارقة «العاجل»', 'size': 15, 'color': RED, 'bold': True, 'space': 4},
     {'t': 'من 13 طلباً «عاجلاً»: تأخّر 8 وأُلغي 1 — أي أن صفة العاجل لا تضمن السرعة.', 'size': 14, 'color': INK}],
    align=PP_ALIGN.RIGHT)
decision(s, '12 طلباً لم تكتمل (ملغاة/قيد المراجعة)، وثلثا الأوامر متأخرة — مؤشر على خلل في الدورة.')
source_note(s)
page(s)

# لوحة 3) أين يضيع الوقت
s = new(); brandbar(s, 'أين يضيع الوقت؟ تشخيص التأخير', 'سبب التأخير الأبرز')
big = [
    ('44٪', 'من الطلبات المتأخرة', 'سببها التحويل المالي', RED),
    ('336', 'يوم تأخير', 'من أصل 688 بسبب المالية', RED),
    ('15.6', 'يوم متوسط', 'زمن التوريد الكلي', AMBER),
]
cw = (SW - Inches(1.2)) / 3
for i, (b, l1, l2, co) in enumerate(big):
    x = SW - Inches(0.5) - (i + 1) * cw - i * Inches(0.1) + Inches(0.1)
    kpi_tile(s, x, Inches(1.55), cw, Inches(1.9), b, l1, co, l2)
txt(s, SW - Inches(6.6), Inches(3.75), Inches(6.3), Inches(0.4),
    'أوامر الشراء حسب القطاع', size=16, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)
hbars(s, Inches(0.5), Inches(4.25), Inches(12.3), [
    ('الصيانة والتشغيل', 27, INK), ('الإنشاءات', 12, AMBER),
    ('الإدارة العامة', 9, GREEN), ('النقليات', 4, GOLD),
], maxv=27, lblw=Inches(3.0), bar_h=Inches(0.34), gap=Inches(0.16))
decision(s, 'نحو نصف أيام التأخير سببها التحويل المالي — ما يؤكد أولوية إلزام المالية بمدة زمنية.')
source_note(s)
page(s)

# 5–11) المحاور السبعة (عمودان)
axes = [
    ('أولاً: تنظيم الطلبات', [
        ('طلبات بالبريد والواتساب بصفة «عاجل» بلا نموذج', 'نموذج طلب رسمي موحّد موقّع، ولا يُعتدّ بغيره'),
        ('الإفراط في «عاجل» وأغلبه لا يُنفّذ', '«عاجل» باعتماد خطّي مبرّر وسقف شهري'),
        ('طلبات معلّقة دون إجراء (خاصة لدى المالية)', 'مدة بتّ محددة، وتصعيد متدرّج: مدير الإدارة ثم الإدارة العليا عند التجاوز'),
        ('غموض الدورة ومن يعتمد', 'دورة موثّقة ومصفوفة صلاحيات'),
        ('مواصفات ناقصة وقنوات متفرّقة', 'إرفاق صورة/علامة تجارية، وقناة رسمية واحدة'),
    ]),
    ('ثانياً: ازدواجية التوريد والشراء المباشر', [
        ('المشتريات تُورّد صنفاً ورّدته الصيانة', 'إشعار فوري بالاكتفاء، ومنع توريد موازٍ بعد أمر الشراء'),
        ('شراء مباشر خارج المشتريات', 'كل شراء وتوريد عبر المشتريات حصراً'),
        ('شراء من موردين بعينهم دون رجوع', 'إدراجهم في قائمة معتمدة تُدار مركزياً'),
        ('شراء صنف متوفّر بالمخزون', 'فحص المستودع إلزامي قبل الطلب'),
    ]),
    ('ثالثاً: تسعير المناقصات وسجل الأسعار', [
        ('إعادة تجميع الأسعار يدوياً لكل مناقصة', 'الرجوع لسجل الأسعار المرجعي المُفعّل في النظام'),
        ('السجل يحتاج تغذية مستمرة بالأسعار', 'مؤتمت داخل النظام؛ ويُثرى بربط فواتير النظام المالي'),
        ('ضيق وقت تسعير المناقصات', 'تخصيص مدة معيارية كافية متّفق عليها'),
    ]),
    ('رابعاً: الموردون والعلاقات', [
        ('تآكل ثقة الموردين بطلباتنا', 'ضبط «العاجل» والالتزام بإغلاق الطلبات'),
        ('طلبات شهرية متكررة بلا عقود', 'عقود إطارية سنوية مُرساة بثلاثة عروض'),
        ('الطلبات المستديمة المتكررة', 'تخضع لثلاثة عروض كاملة كغيرها دون استثناء'),
        ('مشاكل الصيانة والتشغيل مع الموردين غير محصورة', 'تسليم المشتريات بها وبالمسؤول عنها لتوحيد التعامل'),
        ('غموض موردي الصيانة والتشغيل والإدارات (نقداً/أجلاً)', 'حصرهم وتسجيلهم لدينا وحفظ سجلاتهم في قائمة موحّدة'),
    ]),
    ('خامساً: العلاقة مع الصيانة والتشغيل', [
        ('لا إطار منظّم للعلاقة', 'اتفاقية تحدّد المسارات الأربعة'),
        ('لا جهة تنسيق واحدة', 'منسّق من الصيانة ملمّ بالمشتريات'),
        ('آليات مختلفة لكل طلب', 'آلية موحّدة لكل نوع طلب'),
    ]),
    ('سادساً: العهد والموارد', [
        ('عهدة قصر الياسمين مسجّلة على مندوب المشتريات رغم أن الشراء والصرف يتمّ عبر المشروع', 'نقل العهدة إلى مسؤول المشروع باعتباره المنفّذ الفعلي'),
        ('بقاء العهدة على المندوب يحمّله مسؤولية مالية عن مصروفات لا يديرها', 'تسوية العهدة الحالية وإبراء ذمّة المندوب فوراً'),
        ('لا عهدة نقدية للمشتريات', 'عهدة مستديمة بضوابط صرف'),
        ('لا وسيلة توريد للمواقع', 'سائق تابع للمشتريات'),
    ]),
    ('سابعاً: التكامل مع الإدارة المالية', [
        ('بطء التحويل للموردين وتعليق المعاملات', 'مدة ملزمة (3–5 أيام)، وتصعيد للإدارة العليا عند التجاوز'),
        ('إثراء سجل الأسعار بأسعار الشراء الفعلية', 'ربط فواتير النظام المالي بسجل الأسعار المُفعّل'),
        ('لا مسار طوارئ محكوم', 'مسار بموافقة مسبقة وسقف محدّد'),
    ]),
]
def price_register_slide():
    s = new(); brandbar(s, 'إنجاز سجل الأسعار', 'منظومة فعلية مبنية ومُشغّلة داخل نظام المشتريات')
    tiles = [
        (PR_ITEMS, 'صنف مسجّل', INK, '579 مُسعّراً · 6 بلا سعر'),
        (PR_PRICES, 'سجل سعري', GOLD, 'تاريخ أسعار مرجعي للمناقصات'),
        ('25', 'فئة', GREEN, 'تغطّي مواد المشاريع'),
        ('53', 'مورد معتمد', INK, 'شركاء نشطون بسجلّاتهم'),
    ]
    cw = (SW - Inches(1.2)) / 4
    for i, (big, lbl, co, sub) in enumerate(tiles):
        x = SW - Inches(0.5) - (i + 1) * cw - i * Inches(0.07) + Inches(0.07)
        kpi_tile(s, x, Inches(1.5), cw, Inches(1.9), big, lbl, co, sub)
    rect(s, Inches(0.5), Inches(3.7), SW - Inches(1.0), Inches(1.55), CARD, line=LINE)
    rect(s, SW - Inches(0.5) - Pt(8), Inches(3.7), Pt(8), Inches(1.55), GOLD)
    txt(s, Inches(0.7), Inches(3.85), SW - Inches(1.4), Inches(1.3),
        [{'t': 'من الصفر إلى منظومة مؤتمتة', 'size': 16, 'color': GOLD, 'bold': True, 'space': 5},
         {'t': 'عند استلام الإدارة لم يكن هناك سجل أسعار، وكان عدد الموردين المستلَمين 4 فقط؛', 'size': 15, 'color': INK, 'bold': True, 'space': 3},
         {'t': 'واليوم: نظام مشتريات فعّال ومؤتمت، و53 مورداً معتمداً في 25 فئة، و585 صنفاً و864 سجلاً سعرياً يُحدَّث باستمرار.', 'size': 15, 'color': SLATE}],
        align=PP_ALIGN.RIGHT)
    source_note(s, 'المصدر: نظام المشتريات — لوحة التحكم (سجل الأصناف والأسعار والموردين)')
    page(s)


for k, (title, rows) in enumerate(axes):
    s = new(); brandbar(s, title, P1)
    two_col(s, Inches(1.5), rows, fsize=15, row_h=0.78)
    page(s)
    if k == 2:   # بعد محور تسعير المناقصات وسجل الأسعار
        price_register_slide()

# 12) السياسات الحاكمة
s = new(); brandbar(s, 'سياسات المشتريات الحاكمة', 'مقترح دليل السياسات والإجراءات — الإصدار 1.1')
pols = [
    ('القناة الموحّدة', 'كل شراء عبر المشتريات حصراً'),
    ('ثلاثة عروض إلزامية', 'لأي شراء مهما كان حجمه'),
    ('مصفوفة الصلاحيات', 'الاعتماد حسب القيمة'),
    ('الدورة المستندية', 'مسار موحّد من الاحتياج للإغلاق'),
    ('اعتماد الموردين', 'قائمة معتمدة وتقييم دوري'),
    ('العقود الإطارية', 'للطلبات المتكررة والمستديمة'),
    ('سجل أسعار المناقصات', 'للمناقصات فقط لا للشراء'),
    ('مسار الطوارئ', 'بموافقة مسبقة موثّقة'),
    ('عهدة المشتريات', 'عهدة مستديمة بضوابط'),
    ('لا دفع قبل أمر الشراء', 'قنوات الدفع: تحويل أو نقاط بيع'),
    ('الفصل بين المهام', 'فصل الطلب والتقييم والاعتماد والدفع'),
    ('التحقق قبل الاعتماد', 'طلب وأمر وعروض ومقارنة'),
]
cw = (SW - Inches(1.1)) / 3
for i, (ti, de) in enumerate(pols):
    r, c = divmod(i, 3)
    x = SW - Inches(0.45) - (c + 1) * cw - c * Inches(0.05) + Inches(0.05)
    y = Inches(1.55) + r * Inches(1.28)
    rect(s, x, y, cw, Inches(1.12), CARD, line=LINE)
    rect(s, x + cw - Pt(7), y, Pt(7), Inches(1.12), GOLD)
    txt(s, x + Inches(0.12), y + Inches(0.1), cw - Inches(0.3), Inches(0.95),
        [{'t': ti, 'size': 16, 'color': INK, 'bold': True, 'space': 3},
         {'t': de, 'size': 13, 'color': SLATE}], align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE)
page(s, 12)

# 13) الدورة الكاملة (خريطة)
s = new(); brandbar(s, 'الدورة الكاملة للمشتريات', P2)
steps = ['طلب الشراء', 'طلب ثلاثة\nعروض', 'التحليل\nوالمفاضلة', 'تعميد المورد',
         'أمر الشراء', 'الصرف/التحويل', 'التوريد\nوالاستلام', 'المطابقة\nوالإغلاق']
per = 5; bw = Inches(2.32); bh = Inches(1.1)
gapx = (SW - Inches(0.6) - per * bw) / (per - 1)
for i, st in enumerate(steps):
    row = i // per; col = i % per
    x = SW - Inches(0.3) - bw - col * (bw + gapx)
    y = Inches(1.85) + row * Inches(2.05)
    co = GOLD if i in (1, 2, 3) else INK
    tc = INK if co == GOLD else WHITE
    rect(s, x, y, bw, bh, co)
    txt(s, x, y, bw, bh, st.replace('\n', ' '), size=15, color=tc, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if col < per - 1:
        txt(s, x - gapx, y, gapx, bh, '◄', size=22, color=MUTE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, is_rtl=False)
    elif row == 0:
        txt(s, x, y + bh, bw, Inches(0.95), '▼', size=22, color=MUTE,
            align=PP_ALIGN.CENTER, is_rtl=False)
decision(s, 'ثلاثة عروض إلزامية في كل دورة — ولا توريد خارج هذا المسار.')
page(s, 13)

# 14) إجراءات الدورة ومددها (مقترح الدليل 1.1)
s = new(); brandbar(s, 'إجراءات الدورة ومددها', 'مقترح دليل السياسات والإجراءات — الإصدار 1.1')
sop_rows = [
    ('1) طلب الشراء', 'تحديد الحاجة والمواصفات + تحقق مالي مسبق', 'الجهة الطالبة', 'يومان'),
    ('2) طلب العروض', 'دعوة ثلاثة موردين مؤهلين على الأقل', 'المشتريات', '3–5 أيام'),
    ('3) التحليل والمفاضلة', 'مطابقة فنية وجدول مقارنة', 'المشتريات', 'يوم'),
    ('4) تعميد المورد', 'اعتماد الأفضل قيمة بعد ثلاثة عروض', 'مدير العمليات', 'يوم'),
    ('5) أمر الشراء', 'إصدار الأمر وإرساله للمورد', 'المشتريات', 'يوم'),
    ('6) الصرف/التحويل', 'الدفع بعد أمر الشراء (تحويل/نقاط بيع)', 'المالية', '—'),
    ('7) التوريد والاستلام', 'فحص الكمية والجودة ومحضر استلام', 'المستلم/المستودعات', 'يوم'),
    ('8) المطابقة والإغلاق', 'مطابقة المستندات وأرشفة الملف', 'المالية والمشتريات', 'يوم'),
]
g = s.shapes.add_table(len(sop_rows) + 1, 4, Inches(0.4), Inches(1.5),
                       SW - Inches(0.8), Inches(0.6 * (len(sop_rows) + 1)))
tb = g.table; table_rtl(tb)
for i, frac in enumerate((2.6, 5.2, 2.4, 1.6)):
    tb.columns[i].width = Emu(int((SW - Inches(0.8)) * frac / 11.8))
for j, htxt in enumerate(('المرحلة', 'الإجراء', 'المسؤول', 'المدة')):
    c = tb.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = INK
    c.vertical_anchor = MSO_ANCHOR.MIDDLE; c.margin_left = c.margin_right = Pt(5)
    p = c.text_frame.paragraphs[0]; rtl(p); p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = htxt; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = F
palette = [(INK, True, PP_ALIGN.RIGHT), (SLATE, False, PP_ALIGN.RIGHT),
           (INK, False, PP_ALIGN.CENTER), (GREEN, True, PP_ALIGN.CENTER)]
for i, row in enumerate(sop_rows, start=1):
    for j, cell in enumerate(row):
        c = tb.cell(i, j); c.fill.solid()
        c.fill.fore_color.rgb = CARD if i % 2 else ROWALT
        c.vertical_anchor = MSO_ANCHOR.MIDDLE; c.margin_left = c.margin_right = Pt(6)
        col, bd, al = palette[j]
        p = c.text_frame.paragraphs[0]; rtl(p); p.alignment = al
        r = p.add_run(); r.text = cell; r.font.size = Pt(12.5); r.font.name = F
        r.font.color.rgb = col; r.font.bold = bd
page(s)

# 15) مصفوفة الصلاحيات
s = new(); brandbar(s, 'مصفوفة الصلاحيات', P2)
tiers = [
    ('المستوى الأول', 'حتى 25,000', 'مدير المشتريات', GREEN),
    ('المستوى الثاني', 'حتى 100,000', 'المشتريات والمالية', AMBER),
    ('المستوى الثالث', 'حتى 500,000', 'المدير العام', RED),
    ('المستوى الرابع', 'أكثر من 500,000', 'المدير العام ولجنة', INK),
]
cw = (SW - Inches(1.1)) / 4
for i, (ti, val, who, co) in enumerate(tiers):
    x = SW - Inches(0.45) - (i + 1) * cw - i * Inches(0.05) + Inches(0.05); y = Inches(1.7)
    rect(s, x, y, cw, Inches(0.62), co)
    txt(s, x, y, cw, Inches(0.62), ti, size=16, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, x, y + Inches(0.68), cw, Inches(1.7), CARD, line=LINE)
    txt(s, x + Inches(0.08), y + Inches(0.82), cw - Inches(0.16), Inches(1.5),
        [{'t': val + ' ريال', 'size': 18, 'color': INK, 'bold': True, 'space': 10},
         {'t': who, 'size': 14, 'color': SLATE}], align=PP_ALIGN.CENTER)
rect(s, Inches(0.5), Inches(4.45), SW - Inches(1.0), Inches(0.62), GOLD)
txt(s, Inches(0.5), Inches(4.45), SW - Inches(1.0), Inches(0.62),
    'العروض الثلاثة قاعدة ثابتة في كل المستويات', size=17, color=INK, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.5), Inches(5.2), SW - Inches(1.0), Inches(0.7),
    'الاستثناءات (باعتماد مدير المشتريات): النثريات — المتعاقد عليه إطارياً — الطوارئ المعتمدة مسبقاً.',
    size=14, color=SLATE, align=PP_ALIGN.CENTER)
page(s)

# المسارات الخاصة والرقابة (من الدليل)
s = new(); brandbar(s, 'المسارات الخاصة والرقابة', 'مقترح دليل السياسات والإجراءات — الإصدار 1.1')
ctrl = [
    ('المشتريات الطارئة', 'بإذن مدير العمليات، واستكمال المنافسة والتبرير لاحقاً'),
    ('وحيد المصدر', 'بنموذج تبرير فني/تجاري معتمد عند تعذّر المنافسة'),
    ('المشتريات المتكررة', 'عقود إطارية سنوية مع موردين معتمدين'),
    ('الفصل بين الواجبات', 'فصل الطلب والتقييم والاعتماد والدفع بين أطراف'),
    ('حفظ السجلات', 'لا تقل عن خمس سنوات لكل ملفات المشتريات'),
    ('التدقيق الدوري', 'تدقيق نصف سنوي على المنافسة والموافقات والأرشفة'),
]
cw = (SW - Inches(1.1)) / 3
for i, (ti, de) in enumerate(ctrl):
    r, c = divmod(i, 3)
    x = SW - Inches(0.45) - (c + 1) * cw - c * Inches(0.05) + Inches(0.05)
    y = Inches(1.7) + r * Inches(1.9)
    rect(s, x, y, cw, Inches(1.65), CARD, line=LINE)
    rect(s, x + cw - Pt(7), y, Pt(7), Inches(1.65), GOLD)
    txt(s, x + Inches(0.14), y + Inches(0.16), cw - Inches(0.34), Inches(1.35),
        [{'t': ti, 'size': 18, 'color': INK, 'bold': True, 'space': 5},
         {'t': de, 'size': 14, 'color': SLATE}], align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE)
page(s)

# 16) خارطة الطريق
s = new(); brandbar(s, 'خارطة الطريق', P3)
phases = [
    ('أول 30 يوماً', GREEN, ['النموذج الموحّد', 'مصفوفة الصلاحيات', 'حصر الشراء بالمشتريات', 'التزام المالية الزمني']),
    ('حتى 90 يوماً', AMBER, ['أتمتة سجل الأسعار', 'العقود الإطارية', 'عهدة وسائق للمشتريات', 'معالجة عهدة قصر الياسمين']),
    ('المدى البعيد', INK, ['نظام مشتريات موحّد', 'لوحات متابعة', 'تقييم الموردين', 'استراتيجية التوريد']),
]
cw = (SW - Inches(1.1)) / 3
for i, (ti, co, items) in enumerate(phases):
    x = SW - Inches(0.45) - (i + 1) * cw - i * Inches(0.05) + Inches(0.05)
    rect(s, x, Inches(1.7), cw, Inches(0.8), co)
    txt(s, x, Inches(1.7), cw, Inches(0.8), ti, size=22, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, x, Inches(2.55), cw, Inches(3.3), CARD, line=LINE)
    bullets(s, x + Inches(0.15), Inches(2.75), cw - Inches(0.35), Inches(3.0), items,
            size=16, gap=16, color=SLATE)
page(s, 16)

# 17) مؤشرات النجاح: الحالي ← الهدف
s = new(); brandbar(s, 'مؤشرات النجاح: الحالي والهدف', P3)
kpis = [
    ('الطلبات المتأخرة', '68٪', 'أقل من 15٪'),
    ('متوسط زمن التوريد', '15.6 يوم', '8 أيام'),
    ('تأخيرات بسبب المالية', '44٪', 'أقل من 10٪'),
    ('الطلبات غير المكتملة', '23٪', 'أقل من 5٪'),
    ('الالتزام بالعروض الثلاثة', 'يُبدأ القياس', '95٪ فأكثر'),
    ('الموردون المسجّلون', '42 مورداً', 'نموّ مستمر'),
]
cw = (SW - Inches(1.2)) / 3
for i, (ti, cur, tg) in enumerate(kpis):
    r, c = divmod(i, 3)
    x = SW - Inches(0.5) - (c + 1) * cw - c * Inches(0.1) + Inches(0.1)
    y = Inches(1.65) + r * Inches(2.15)
    rect(s, x, y, cw, Inches(1.85), CARD, line=LINE)
    rect(s, x, y, cw, Inches(0.55), INK)
    txt(s, x, y, cw, Inches(0.55), ti, size=15, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    half = cw / 2
    txt(s, x + half, y + Inches(0.62), half, Inches(1.15),
        [{'t': 'الحالي', 'size': 11, 'color': MUTE, 'bold': True, 'space': 3},
         {'t': cur, 'size': 19, 'color': RED, 'bold': True}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, y + Inches(0.62), half, Inches(1.15),
        [{'t': 'الهدف', 'size': 11, 'color': MUTE, 'bold': True, 'space': 3},
         {'t': tg, 'size': 19, 'color': GREEN, 'bold': True}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + half - Inches(0.12), y + Inches(0.75), Inches(0.24), Inches(0.9), '←',
        size=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, is_rtl=False)
page(s)

# 18) القرارات المطلوبة
s = new(); brandbar(s, 'القرارات المطلوبة', P3)
asks = [
    'حصر كل شراء عبر المشتريات حصراً.',
    'اعتماد السياسات والدورة الكاملة.',
    'إقرار مصفوفة الصلاحيات والعروض الثلاثة.',
    'إلزام المالية بمدة التحويل.',
    'إنشاء سجل أسعار للمناقصات.',
    'تفويض العقود الإطارية.',
    'معالجة العهد وموارد المشتريات.',
    'اعتماد خارطة الطريق والمؤشرات.',
]
cw = (SW - Inches(1.2)) / 2
for i, a in enumerate(asks):
    r, c = divmod(i, 2)
    x = SW - Inches(0.5) - (c + 1) * cw - c * Inches(0.2) + Inches(0.2)
    y = Inches(1.65) + r * Inches(1.05)
    rect(s, x, y, cw, Inches(0.86), CARD, line=LINE)
    rect(s, x + cw - Inches(0.75), y, Inches(0.75), Inches(0.86), GREEN)
    txt(s, x + cw - Inches(0.75), y, Inches(0.75), Inches(0.86), str(i + 1), size=26,
        color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, is_rtl=False)
    txt(s, x + Inches(0.15), y, cw - Inches(0.95), Inches(0.86), a, size=15,
        color=INK, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
page(s, 18)

# 19) الخاتمة
s = new()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(2.95), SW, Pt(4), GOLD)
txt(s, Inches(1), Inches(1.9), SW - Inches(2), Inches(1.1),
    'منظومة واضحة تحمي المال وتنظّم التوريد.', size=34, color=WHITE, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.35), SW - Inches(2), Inches(0.8),
    'وفريق المشتريات جاهز للتنفيذ فور الاعتماد.', size=18, color=GOLD,
    align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.4), SW - Inches(2), Inches(0.5),
    'مع وافر الاحترام والتقدير — إدارة العمليات والمشتريات', size=15, color=SAND,
    align=PP_ALIGN.CENTER)

TOTAL = len(slides)
out = os.path.join(os.path.dirname(__file__), 'procurement-executive-deck.pptx')
prs.save(out)
print('Saved', out, '|', os.path.getsize(out), 'bytes |', TOTAL, 'slides')
